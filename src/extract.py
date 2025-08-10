import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def extract_text_and_images(pptx_path, output_image_dir):
    prs = Presentation(pptx_path)
    slides_data = []

    os.makedirs(output_image_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        image_files = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                img = Image.open(io.BytesIO(image.blob))

                filename = f"slide{slide_num}_{len(image_files)+1}.png"
                filepath = os.path.join(output_image_dir, filename)
                img.save(filepath)
                image_files.append(filepath)

        slides_data.append({
            "slide_number": slide_num,
            "texts": slide_texts,
            "images": image_files
        })

    return slides_data
